#!/usr/bin/env python3
"""
PDF to Office Converter
Converts PDF files to Word (.docx), Excel (.xlsx), and PowerPoint (.pptx)
"""

import sys
import os
from pathlib import Path

def convert_pdf_to_word(pdf_path: str, output_path: str) -> bool:
    """Convert PDF to Word document using pdf2docx with PyMuPDF fallback"""
    try:
        from pdf2docx import Converter
        
        print("[INFO] Using pdf2docx for conversion...")
        cv = Converter(pdf_path)
        cv.convert(output_path)
        cv.close()
        
        print(f"Successfully converted to Word: {output_path}")
        return True
    except Exception as e:
        print(f"pdf2docx error: {e}", file=sys.stderr)
        print("[INFO] Trying fallback method with PyMuPDF...", file=sys.stderr)
        
        # Fallback to PyMuPDF for text extraction
        try:
            import fitz  # PyMuPDF
            from docx import Document
            from docx.shared import Pt
            
            doc = Document()
            pdf_doc = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                text = page.get_text()
                
                if page_num > 0:
                    doc.add_page_break()
                
                # Add text to document
                if text.strip():
                    paragraph = doc.add_paragraph(text)
                    paragraph.style.font.size = Pt(11)
            
            pdf_doc.close()
            doc.save(output_path)
            
            print(f"Successfully converted to Word using PyMuPDF: {output_path}")
            return True
        except Exception as fallback_error:
            print(f"Fallback conversion also failed: {fallback_error}", file=sys.stderr)
            return False

def convert_pdf_to_excel(pdf_path: str, output_path: str) -> bool:
    """Convert PDF tables to Excel using tabula-py with memory limits"""
    try:
        import tabula
        from openpyxl import Workbook
        import os
        
        # Set Java memory limits to prevent OOM on Railway (512MB max)
        java_options = [
            '-Xmx512m',  # Maximum heap size
            '-Xms128m',  # Initial heap size
        ]
        
        # Read all tables from PDF with memory-constrained Java
        print("[INFO] Extracting tables from PDF...")
        tables = tabula.read_pdf(
            pdf_path, 
            pages='all', 
            multiple_tables=True,
            java_options=java_options
        )
        
        if not tables or len(tables) == 0:
            print("No tables found in PDF", file=sys.stderr)
            return False
        
        print(f"[INFO] Found {len(tables)} table(s)")
        
        # Create Excel workbook
        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet
        
        # Add each table as a separate sheet
        for idx, table in enumerate(tables):
            sheet_name = f"Table_{idx + 1}"
            ws = wb.create_sheet(title=sheet_name)
            
            # Write headers
            for col_idx, col_name in enumerate(table.columns, start=1):
                ws.cell(row=1, column=col_idx, value=str(col_name))
            
            # Write data
            for row_idx, row in enumerate(table.itertuples(index=False), start=2):
                for col_idx, value in enumerate(row, start=1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
        
        wb.save(output_path)
        print(f"Successfully converted to Excel: {output_path}")
        return True
    except Exception as e:
        print(f"Error converting PDF to Excel: {e}", file=sys.stderr)
        return False

def convert_pdf_to_ppt(pdf_path: str, output_path: str) -> bool:
    """Convert PDF pages to PowerPoint slides using pdf2image and python-pptx"""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import tempfile
        
        # Convert PDF pages to images
        images = convert_from_path(pdf_path)
        
        if not images:
            print("No pages found in PDF", file=sys.stderr)
            return False
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add each page as a slide
        for idx, image in enumerate(images):
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Save image temporarily
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                image.save(tmp_path, 'PNG')
            
            try:
                # Add image to slide
                left = Inches(0)
                top = Inches(0)
                slide.shapes.add_picture(tmp_path, left, top, width=prs.slide_width)
            finally:
                # Clean up temporary file
                os.unlink(tmp_path)
        
        prs.save(output_path)
        print(f"Successfully converted to PowerPoint: {output_path}")
        return True
    except Exception as e:
        print(f"Error converting PDF to PowerPoint: {e}", file=sys.stderr)
        return False

def main():
    """Main entry point for PDF conversion"""
    if len(sys.argv) != 4:
        print("Usage: python pdf_converter.py <input.pdf> <output.docx|xlsx|pptx> <format>")
        print("Format: word, excel, or ppt")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    output_format = sys.argv[3].lower()
    
    # Validate input file
    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)
    
    # Validate output format
    if output_format not in ['word', 'excel', 'ppt']:
        print(f"Error: Unsupported format: {output_format}", file=sys.stderr)
        print("Supported formats: word, excel, ppt")
        sys.exit(1)
    
    # Perform conversion
    success = False
    
    if output_format == 'word':
        success = convert_pdf_to_word(input_path, output_path)
    elif output_format == 'excel':
        success = convert_pdf_to_excel(input_path, output_path)
    elif output_format == 'ppt':
        success = convert_pdf_to_ppt(input_path, output_path)
    
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
